"""
文档解析服务
支持Word、PDF、PPT、Markdown、HTML、TXT、CSV文档解析
"""
import logging
import os
from typing import Dict, List, Optional
from pathlib import Path

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logging.warning("python-docx未安装，Word文档解析功能不可用")

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logging.warning("PyPDF2未安装，PDF文档解析功能不可用")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx未安装，PPT文档解析功能不可用")

try:
    import markdown
    MARKDOWN_AVAILABLE = True
except ImportError:
    MARKDOWN_AVAILABLE = False
    logging.warning("python-markdown未安装，Markdown文档解析功能不可用")

try:
    from bs4 import BeautifulSoup
    HTML_AVAILABLE = True
except ImportError:
    HTML_AVAILABLE = False
    logging.warning("beautifulsoup4未安装，HTML文档解析功能不可用")

logger = logging.getLogger(__name__)


class DocumentParserService:
    """文档解析服务"""
    
    def __init__(self):
        self.supported_formats = []
        if DOCX_AVAILABLE:
            self.supported_formats.extend(["docx", "doc"])
        if PDF_AVAILABLE:
            self.supported_formats.append("pdf")
        if PPTX_AVAILABLE:
            self.supported_formats.extend(["pptx", "ppt"])
        if MARKDOWN_AVAILABLE:
            self.supported_formats.extend(["md", "markdown"])
        if HTML_AVAILABLE:
            self.supported_formats.extend(["html", "htm"])
        # TXT和CSV总是支持的
        self.supported_formats.extend(["txt", "csv"])
    
    def parse_document(self, file_path: str, extract_images: bool = False) -> Dict:
        """
        解析文档（增强版）
        
        Args:
            file_path: 文档文件路径
            extract_images: 是否提取图片的OCR文字
            
        Returns:
            解析结果字典，包含：
            - content: 文档文本内容
            - structure: 文档结构信息（标题、段落、表格等）
            - metadata: 文档元数据
            - chunks: 预分块的内容（可选）
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        file_extension = Path(file_path).suffix.lower().lstrip('.')
        
        logger.info(f"开始解析文档: 文件路径={file_path}, 文件类型={file_extension}")
        
        try:
            if file_extension in ["docx", "doc"]:
                return self._parse_word(file_path)
            elif file_extension == "pdf":
                return self._parse_pdf(file_path)
            elif file_extension in ["pptx", "ppt"]:
                return self._parse_pptx(file_path)
            elif file_extension in ["md", "markdown"]:
                return self._parse_markdown(file_path)
            elif file_extension in ["html", "htm"]:
                return self._parse_html(file_path)
            elif file_extension == "txt":
                return self._parse_txt(file_path)
            elif file_extension == "csv":
                return self._parse_csv(file_path)
            else:
                raise ValueError(f"不支持的文件格式: {file_extension}")
        except Exception as e:
            logger.error(f"文档解析失败: 文件路径={file_path}, 错误={str(e)}", exc_info=True)
            raise
    
    def _parse_word(self, file_path: str) -> Dict:
        """解析Word文档"""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx未安装，无法解析Word文档")
        
        try:
            doc = Document(file_path)
            
            # 提取文本内容
            paragraphs = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append(text)
            
            content = "\n".join(paragraphs)
            
            # 提取标题结构
            structure = self._extract_structure(paragraphs)
            
            # 提取元数据
            metadata = {
                "title": doc.core_properties.title or "",
                "author": doc.core_properties.author or "",
                "created": str(doc.core_properties.created) if doc.core_properties.created else "",
                "modified": str(doc.core_properties.modified) if doc.core_properties.modified else "",
            }
            
            logger.info(f"Word文档解析成功: 段落数={len(paragraphs)}, 字符数={len(content)}")
            
            return {
                "content": content,
                "structure": structure,
                "metadata": metadata,
                "paragraph_count": len(paragraphs),
                "char_count": len(content)
            }
            
        except Exception as e:
            logger.error(f"Word文档解析失败: {str(e)}", exc_info=True)
            raise
    
    def _parse_pdf(self, file_path: str) -> Dict:
        """解析PDF文档"""
        if not PDF_AVAILABLE:
            raise ImportError("PyPDF2未安装，无法解析PDF文档")
        
        try:
            paragraphs = []
            metadata = {}
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                # 提取元数据
                if pdf_reader.metadata:
                    metadata = {
                        "title": pdf_reader.metadata.get("/Title", ""),
                        "author": pdf_reader.metadata.get("/Author", ""),
                        "subject": pdf_reader.metadata.get("/Subject", ""),
                        "creator": pdf_reader.metadata.get("/Creator", ""),
                    }
                
                # 提取文本内容
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        text = page.extract_text()
                        if text.strip():
                            paragraphs.append(text.strip())
                    except Exception as e:
                        logger.warning(f"PDF第{page_num + 1}页解析失败: {str(e)}")
                        continue
            
            content = "\n".join(paragraphs)
            
            # 提取结构
            structure = self._extract_structure(paragraphs)
            
            logger.info(f"PDF文档解析成功: 页数={len(pdf_reader.pages)}, 段落数={len(paragraphs)}, 字符数={len(content)}")
            
            return {
                "content": content,
                "structure": structure,
                "metadata": metadata,
                "page_count": len(pdf_reader.pages) if 'pdf_reader' in locals() else 0,
                "paragraph_count": len(paragraphs),
                "char_count": len(content)
            }
            
        except Exception as e:
            logger.error(f"PDF文档解析失败: {str(e)}", exc_info=True)
            raise
    
    def _parse_pptx(self, file_path: str) -> Dict:
        """解析PPTX文档"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx未安装，无法解析PPTX文档")
        
        try:
            prs = Presentation(file_path)
            
            # 提取幻灯片内容
            slides_content = []
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                
                if slide_text:
                    slides_content.append("\n".join(slide_text))
            
            content = "\n\n".join(slides_content)
            
            # 提取元数据
            metadata = {
                "title": Path(file_path).stem,
                "slide_count": len(prs.slides),
                "file_size": os.path.getsize(file_path),
                "created_time": self._get_file_created_time(file_path),
                "modified_time": self._get_file_modified_time(file_path)
            }
            
            # 提取结构
            structure = {
                "title": metadata["title"],
                "slide_count": len(prs.slides),
                "sections": []
            }
            
            logger.info(f"PPTX文档解析成功: 幻灯片数={len(prs.slides)}, 字符数={len(content)}")
            
            return {
                "content": content,
                "structure": structure,
                "metadata": metadata,
                "slide_count": len(prs.slides),
                "char_count": len(content)
            }
            
        except Exception as e:
            logger.error(f"PPTX文档解析失败: {str(e)}", exc_info=True)
            raise
    
    def _parse_markdown(self, file_path: str) -> Dict:
        """解析Markdown文档"""
        if not MARKDOWN_AVAILABLE:
            raise ImportError("python-markdown未安装，无法解析Markdown文档")
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                md_content = f.read()
            
            # 提取标题结构
            structure = self._extract_markdown_structure(md_content)
            
            # 提取元数据
            metadata = {
                "title": structure.get("title", Path(file_path).stem),
                "file_size": os.path.getsize(file_path),
                "created_time": self._get_file_created_time(file_path),
                "modified_time": self._get_file_modified_time(file_path),
                "language": self._detect_language(md_content)
            }
            
            logger.info(f"Markdown文档解析成功: 字符数={len(md_content)}")
            
            return {
                "content": md_content,
                "structure": structure,
                "metadata": metadata,
                "char_count": len(md_content)
            }
            
        except Exception as e:
            logger.error(f"Markdown文档解析失败: {str(e)}", exc_info=True)
            raise
    
    def _parse_html(self, file_path: str) -> Dict:
        """解析HTML文档"""
        if not HTML_AVAILABLE:
            raise ImportError("beautifulsoup4未安装，无法解析HTML文档")
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # 使用BeautifulSoup解析
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # 提取纯文本
            text = soup.get_text(separator='\n', strip=True)
            
            # 提取标题结构
            structure = {
                "title": soup.title.string if soup.title else Path(file_path).stem,
                "headings": []
            }
            
            # 提取所有标题
            for heading in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
                structure["headings"].append({
                    "level": heading.name,
                    "text": heading.get_text(strip=True)
                })
            
            # 提取表格
            tables = self._extract_tables_from_html(soup)
            
            # 提取元数据
            metadata = {
                "title": structure["title"],
                "file_size": os.path.getsize(file_path),
                "created_time": self._get_file_created_time(file_path),
                "modified_time": self._get_file_modified_time(file_path),
                "language": self._detect_language(text),
                "table_count": len(tables)
            }
            
            logger.info(f"HTML文档解析成功: 字符数={len(text)}, 表格数={len(tables)}")
            
            return {
                "content": text,
                "structure": structure,
                "metadata": metadata,
                "tables": tables,
                "char_count": len(text)
            }
            
        except Exception as e:
            logger.error(f"HTML文档解析失败: {str(e)}", exc_info=True)
            raise
    
    def _parse_txt(self, file_path: str) -> Dict:
        """解析TXT文档"""
        try:
            # 尝试多种编码
            encodings = ['utf-8', 'gbk', 'gb2312', 'big5']
            content = None
            used_encoding = None
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    used_encoding = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            if content is None:
                raise ValueError("无法识别文件编码")
            
            # 提取结构
            structure = self._extract_structure(content.split('\n'))
            
            # 提取元数据
            metadata = {
                "title": Path(file_path).stem,
                "file_size": os.path.getsize(file_path),
                "created_time": self._get_file_created_time(file_path),
                "modified_time": self._get_file_modified_time(file_path),
                "encoding": used_encoding,
                "language": self._detect_language(content)
            }
            
            logger.info(f"TXT文档解析成功: 编码={used_encoding}, 字符数={len(content)}")
            
            return {
                "content": content,
                "structure": structure,
                "metadata": metadata,
                "char_count": len(content)
            }
            
        except Exception as e:
            logger.error(f"TXT文档解析失败: {str(e)}", exc_info=True)
            raise
    
    def _parse_csv(self, file_path: str) -> Dict:
        """解析CSV文档"""
        try:
            import csv
            
            rows = []
            with open(file_path, 'r', encoding='utf-8') as f:
                csv_reader = csv.reader(f)
                for row in csv_reader:
                    rows.append(row)
            
            # 将CSV转换为文本
            content = "\n".join([",".join(row) for row in rows])
            
            # 提取元数据
            metadata = {
                "title": Path(file_path).stem,
                "file_size": os.path.getsize(file_path),
                "created_time": self._get_file_created_time(file_path),
                "modified_time": self._get_file_modified_time(file_path),
                "row_count": len(rows),
                "column_count": len(rows[0]) if rows else 0,
                "headers": rows[0] if rows else []
            }
            
            # 提取结构
            structure = {
                "title": metadata["title"],
                "headers": metadata["headers"],
                "row_count": metadata["row_count"],
                "column_count": metadata["column_count"]
            }
            
            logger.info(f"CSV文档解析成功: 行数={len(rows)}, 列数={metadata['column_count']}")
            
            return {
                "content": content,
                "structure": structure,
                "metadata": metadata,
                "row_count": len(rows)
            }
            
        except Exception as e:
            logger.error(f"CSV文档解析失败: {str(e)}", exc_info=True)
            raise
    
    def _extract_structure(self, paragraphs: List[str]) -> Dict:
        """
        提取文档结构
        
        识别标题、章节等结构信息
        """
        structure = {
            "title": "",
            "sections": [],
            "headings": []
        }
        
        if not paragraphs:
            return structure
        
        # 第一段通常作为标题
        if paragraphs:
            structure["title"] = paragraphs[0][:100]  # 取前100个字符作为标题
        
        # 识别标题（简单规则：短段落且包含数字编号的可能是标题）
        headings = []
        for para in paragraphs:
            # 简单的标题识别规则
            if len(para) < 100 and (para.startswith(("第", "一、", "二、", "三、", "1.", "2.", "3.")) or 
                                   para.endswith(("：", ":"))):
                headings.append(para)
        
        structure["headings"] = headings[:20]  # 最多保留20个标题
        
        return structure
    
    def extract_key_info(self, content: str) -> Dict:
        """
        提取关键信息
        
        从文档内容中提取需求相关的关键信息
        """
        # 简单的关键词提取（后续可以用NLP技术优化）
        keywords = []
        
        # 提取可能的关键词（长度2-6的中文词汇）
        import re
        words = re.findall(r'[\u4e00-\u9fa5]{2,6}', content)
        
        # 统计词频
        from collections import Counter
        word_freq = Counter(words)
        
        # 取频率最高的前20个词作为关键词
        keywords = [word for word, freq in word_freq.most_common(20)]
        
        return {
            "keywords": keywords,
            "content_length": len(content),
            "sentence_count": len(re.split(r'[。！？\n]', content))
        }
    
    def _extract_markdown_structure(self, content: str) -> Dict:
        """提取Markdown文档结构"""
        structure = {
            "title": "",
            "sections": [],
            "headings": [],
            "code_blocks": 0,
            "links": 0
        }
        
        lines = content.split('\n')
        
        # 提取标题
        headings = []
        for line in lines:
            if line.startswith('#'):
                level = 0
                for char in line:
                    if char == '#':
                        level += 1
                    else:
                        break
                heading_text = line[level:].strip()
                if heading_text:
                    headings.append({
                        "level": level,
                        "text": heading_text
                    })
                    if not structure["title"] and level == 1:
                        structure["title"] = heading_text
            # 统计代码块
            if line.strip().startswith('```'):
                structure["code_blocks"] += 1
            # 统计链接
            if '](' in line:
                structure["links"] += len(line.split('](')) - 1
        
        structure["headings"] = headings[:20]  # 最多保留20个标题
        
        return structure
    
    def _extract_tables_from_html(self, soup) -> List[Dict]:
        """从HTML中提取表格"""
        tables = []
        
        for table in soup.find_all('table'):
            table_data = []
            rows = table.find_all('tr')
            
            for row in rows:
                cells = []
                # 获取td或th单元格
                for cell in row.find_all(['td', 'th']):
                    cell_text = cell.get_text(strip=True)
                    cells.append(cell_text)
                
                if cells:
                    table_data.append(cells)
            
            if table_data:
                tables.append({
                    "headers": table_data[0] if table_data else [],
                    "rows": table_data[1:] if len(table_data) > 1 else [],
                    "row_count": len(table_data),
                    "column_count": len(table_data[0]) if table_data else 0
                })
        
        return tables
    
    def _get_file_created_time(self, file_path: str) -> str:
        """获取文件创建时间"""
        try:
            import stat
            stat_info = os.stat(file_path)
            if hasattr(stat_info, 'st_ctime'):
                import datetime
                return datetime.datetime.fromtimestamp(stat_info.st_ctime).isoformat()
        except:
            pass
        return ""
    
    def _get_file_modified_time(self, file_path: str) -> str:
        """获取文件修改时间"""
        try:
            stat_info = os.stat(file_path)
            import datetime
            return datetime.datetime.fromtimestamp(stat_info.st_mtime).isoformat()
        except:
            return ""
    
    def _detect_language(self, text: str) -> str:
        """检测文档语言（简单实现）"""
        if not text:
            return "unknown"
        
        # 简单的中英文检测
        chinese_chars = len([c for c in text if '\u4e00' <= c <= '\u9fff'])
        english_chars = len([c for c in text if c.isalpha() and c.isascii()])
        
        total_chars = chinese_chars + english_chars
        if total_chars == 0:
            return "unknown"
        
        chinese_ratio = chinese_chars / total_chars
        
        if chinese_ratio > 0.6:
            return "zh"
        elif chinese_ratio < 0.4:
            return "en"
        else:
            return "mixed"
    
    def extract_tables(self, content: str) -> List[Dict]:
        """从内容中提取表格（通用方法）"""
        tables = []
        
        # 尝试检测Markdown表格
        lines = content.split('\n')
        i = 0
        while i < len(lines):
            line = lines[i].strip()
            
            # 检测Markdown表格（|分隔符）
            if line.startswith('|') and line.endswith('|'):
                headers = [cell.strip() for cell in line.split('|')[1:-1]]
                
                # 跳过分隔线
                i += 1
                if i < len(lines) and lines[i].strip().startswith('|'):
                    separator = lines[i].strip()
                    i += 1
                
                # 收集表格行
                rows = []
                while i < len(lines):
                    row_line = lines[i].strip()
                    if row_line.startswith('|') and row_line.endswith('|'):
                        row_cells = [cell.strip() for cell in row_line.split('|')[1:-1]]
                        rows.append(row_cells)
                        i += 1
                    else:
                        break
                
                if rows:
                    tables.append({
                        "headers": headers,
                        "rows": rows,
                        "row_count": len(rows),
                        "column_count": len(headers)
                    })
            
            i += 1
        
        return tables

